import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt


EXPORT_DIR = "exports"
os.makedirs(EXPORT_DIR, exist_ok=True)


def _render_summary_pdf(pdf, content):
    """Render a structured summary analysis into the PDF."""
    if isinstance(content, str):
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5.5, _sanitize_latin1(content))
        return

    title = content.get("title")
    intro = content.get("introduction") or content.get("summary") or content.get("text")
    points = content.get("points", [])
    conclusion = content.get("conclusion")

    if title:
        pdf.set_font("Helvetica", "B", 12)
        pdf.multi_cell(0, 6, _sanitize_latin1(title))
        pdf.ln(3)

    if intro:
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5.5, _sanitize_latin1(intro))
        pdf.ln(3)

    for point in points:
        text = point if isinstance(point, str) else (point.get("text") or str(point))
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(6, 5.5, _sanitize_latin1("-"))
        pdf.multi_cell(0, 5.5, _sanitize_latin1(text))
        pdf.ln(1)

    if conclusion:
        pdf.ln(2)
        pdf.set_font("Helvetica", "I", 10)
        pdf.set_text_color(80, 80, 80)
        pdf.multi_cell(0, 5.5, _sanitize_latin1(conclusion))
        pdf.set_text_color(0, 0, 0)


def _render_keypoints_pdf(pdf, content):
    """Render structured keypoints analysis into the PDF."""
    if isinstance(content, str):
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5.5, _sanitize_latin1(content))
        return

    keypoints = content.get("keypoints") or content.get("items") or (content if isinstance(content, list) else [])

    for kp in keypoints:
        if isinstance(kp, str):
            pdf.set_font("Helvetica", "", 10)
            pdf.cell(6, 5.5, _sanitize_latin1("-"))
            pdf.multi_cell(0, 5.5, _sanitize_latin1(kp))
            pdf.ln(1)
            continue

        theme = kp.get("theme", "")
        importance = kp.get("importance", "")
        points = kp.get("points", [])
        quote = kp.get("verbatim_quote")

        # Theme header
        if theme:
            label = theme
            if importance:
                imp_labels = {"critical": "Essentiel", "high": "Important", "medium": "Notable", "low": "Mineur"}
                label += f"  ({imp_labels.get(importance, importance)})"
            pdf.set_font("Helvetica", "B", 10)
            pdf.multi_cell(0, 6, _sanitize_latin1(label))
            pdf.ln(1)

        # Points
        for pt in points:
            text = pt if isinstance(pt, str) else (pt.get("text") or str(pt))
            pdf.set_font("Helvetica", "", 10)
            pdf.cell(6, 5.5, _sanitize_latin1("-"))
            pdf.multi_cell(0, 5.5, _sanitize_latin1(text))
            pdf.ln(0.5)

        # Quote
        if quote:
            pdf.set_font("Helvetica", "I", 9)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(6, 5, "")
            pdf.multi_cell(0, 4.5, _sanitize_latin1(f'"{quote}"'))
            pdf.set_text_color(0, 0, 0)

        pdf.ln(3)


def _render_generic_analysis_pdf(pdf, content):
    """Fallback: render analysis content as readable text, not raw JSON."""
    if isinstance(content, str):
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5.5, _sanitize_latin1(content))
        return

    if isinstance(content, dict):
        # Try to extract text-like fields
        text_fields = ["text", "content", "summary", "introduction", "conclusion", "description"]
        for field in text_fields:
            if field in content and isinstance(content[field], str):
                pdf.set_font("Helvetica", "", 10)
                pdf.multi_cell(0, 5.5, _sanitize_latin1(content[field]))
                pdf.ln(2)

        # Render list-like fields
        list_fields = ["points", "items", "questions", "actions", "flashcards", "cards",
                        "slides", "bullets", "faq", "decisions"]
        for field in list_fields:
            if field in content and isinstance(content[field], list):
                for item in content[field]:
                    pdf.set_x(pdf.l_margin)  # Reset cursor to left margin
                    if isinstance(item, str):
                        pdf.set_font("Helvetica", "", 10)
                        pdf.multi_cell(0, 5.5, _sanitize_latin1(f"- {item}"))
                    elif isinstance(item, dict):
                        # Try common sub-fields
                        line = item.get("text") or item.get("question") or item.get("title") or item.get("term") or ""
                        if line:
                            pdf.set_font("Helvetica", "B", 10)
                            pdf.set_x(pdf.l_margin)
                            pdf.multi_cell(0, 5.5, _sanitize_latin1(str(line)))
                        # Sub-detail
                        detail = item.get("answer") or item.get("definition") or item.get("description") or ""
                        if detail:
                            pdf.set_font("Helvetica", "", 9)
                            pdf.set_text_color(60, 60, 60)
                            pdf.set_x(pdf.l_margin)
                            pdf.multi_cell(0, 5, _sanitize_latin1(str(detail)))
                            pdf.set_text_color(0, 0, 0)
                        # Quiz: choices + explanation
                        choices = item.get("choices", [])
                        if choices:
                            letters = "ABCDEFGH"
                            for ci, choice in enumerate(choices):
                                letter = letters[ci] if ci < len(letters) else str(ci+1)
                                correct = item.get("answer", "") == letter
                                pdf.set_font("Helvetica", "B" if correct else "", 9)
                                prefix = ">> " if correct else "   "
                                pdf.set_x(pdf.l_margin)
                                pdf.multi_cell(0, 4.5, _sanitize_latin1(f"{prefix}{letter}. {choice}"))
                            explanation = item.get("explanation", "")
                            if explanation:
                                pdf.set_font("Helvetica", "I", 8)
                                pdf.set_text_color(100, 100, 100)
                                pdf.set_x(pdf.l_margin)
                                pdf.multi_cell(0, 4.5, _sanitize_latin1(explanation))
                                pdf.set_text_color(0, 0, 0)
                    pdf.ln(2)
                pdf.ln(2)
        return

    # Absolute fallback
    pdf.set_font("Helvetica", "", 9)
    pdf.multi_cell(0, 5, _sanitize_latin1(json.dumps(content, indent=2, ensure_ascii=False)))


_ANALYSIS_LABELS = {
    "summary": "Resume",
    "keypoints": "Points cles",
    "actions": "Plan d'actions",
    "quiz": "Quiz",
    "flashcards": "Fiches de revision",
    "mindmap": "Carte mentale",
    "slides": "Diapositives",
    "infographic": "Infographie",
    "tables": "Tableaux",
}


def export_to_pdf(transcription, analyses, output_path: str):
    """Export transcription + analyses to PDF via fpdf2."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    # Use built-in fonts with latin-1 fallback for accented chars
    pdf.set_font("Helvetica", "B", 18)
    title = transcription.filename or "Transcription"
    pdf.cell(0, 12, _sanitize_latin1(title), new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    # Metadata line
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(120, 120, 120)
    lang = transcription.language or "N/A"
    dur = f"{transcription.duration or 0:.0f}s" if transcription.duration else "N/A"
    pdf.cell(0, 6, _sanitize_latin1(f"Langue : {lang}  |  Duree : {dur}"), new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    # Transcription text — render paragraphs with spacing
    pdf.set_font("Helvetica", "B", 13)
    pdf.set_text_color(0, 0, 0)
    pdf.cell(0, 8, "Transcription", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font("Helvetica", "", 10)
    paragraphs = (transcription.text or "").split("\n\n")
    for i, para in enumerate(paragraphs):
        para = para.strip()
        if not para:
            continue
        pdf.multi_cell(0, 5.5, _sanitize_latin1(para))
        if i < len(paragraphs) - 1:
            pdf.ln(3)
    pdf.ln(6)

    # Analyses — formatted by type
    for a in analyses:
        pdf.add_page()
        label = _ANALYSIS_LABELS.get(a.type, a.type.replace("_", " ").title())
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(60, 60, 140)
        pdf.cell(0, 9, _sanitize_latin1(label), new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(0, 0, 0)
        pdf.ln(3)

        # Normalize content: handle JSON strings, double-encoded JSON, etc.
        content = a.content
        if isinstance(content, str):
            try:
                content = json.loads(content)
            except (json.JSONDecodeError, TypeError):
                pass

        if a.type == "summary":
            _render_summary_pdf(pdf, content)
        elif a.type == "keypoints":
            _render_keypoints_pdf(pdf, content)
        else:
            _render_generic_analysis_pdf(pdf, content)

    pdf.output(output_path)


def _sanitize_latin1(text: str) -> str:
    """Replace characters that can't be encoded in latin-1 for fpdf built-in fonts."""
    replacements = {
        "\u2019": "'", "\u2018": "'", "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "-", "\u2026": "...", "\u00a0": " ",
        "\u2022": "-", "\u25cf": "-", "\u2023": ">",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text.encode("latin-1", errors="replace").decode("latin-1")


def export_to_srt(transcription, output_path: str):
    """Export transcription segments to SRT subtitle format."""
    lines = []
    for i, seg in enumerate(transcription.segments or [], 1):
        start = _format_srt_time(seg["start"])
        end = _format_srt_time(seg["end"])
        lines.append(f"{i}")
        lines.append(f"{start} --> {end}")
        lines.append(seg["text"])
        lines.append("")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def export_to_vtt(transcription, output_path: str):
    """Export transcription segments to WebVTT format."""
    lines = ["WEBVTT", ""]
    for seg in transcription.segments or []:
        start = _format_vtt_time(seg["start"])
        end = _format_vtt_time(seg["end"])
        lines.append(f"{start} --> {end}")
        lines.append(seg["text"])
        lines.append("")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def export_to_txt(transcription, output_path: str):
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(transcription.text)


def export_to_json(transcription, analyses, output_path: str):
    data = {
        "id": transcription.id,
        "filename": transcription.filename,
        "language": transcription.language,
        "duration": transcription.duration,
        "text": transcription.text,
        "segments": transcription.segments,
        "analyses": {a.type: a.content for a in analyses},
    }
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def export_to_md(transcription, analyses, output_path: str):
    lines = [
        f"# {transcription.filename}",
        f"**Language:** {transcription.language or 'N/A'} | **Duration:** {transcription.duration or 0:.0f}s",
        "", "## Transcription", "", transcription.text, "",
    ]
    for a in analyses:
        lines.append(f"## {a.type.title()}")
        lines.append(f"```json\n{json.dumps(a.content, indent=2, ensure_ascii=False)}\n```")
        lines.append("")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def export_to_pptx(analyses, output_path: str):
    """Export slides analysis to PowerPoint."""
    prs = Presentation()
    slides_data = None
    for a in analyses:
        if a.type == "slides":
            slides_data = a.content
            break
    if not slides_data or "slides" not in slides_data:
        return
    for slide_data in slides_data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")
        body = slide.placeholders[1]
        body.text = "\n".join(slide_data.get("bullets", []))
    prs.save(output_path)


def _format_srt_time(seconds: float) -> str:
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    ms = int((seconds % 1) * 1000)
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def _format_vtt_time(seconds: float) -> str:
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    ms = int((seconds % 1) * 1000)
    return f"{h:02d}:{m:02d}:{s:02d}.{ms:03d}"
